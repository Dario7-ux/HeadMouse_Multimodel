"""
FocuzVoz 3.0 — Generador de la Infografía Arquitectónica Definitiva
- Diseño B&W industrial premium de alta resolución (28x16 pulgadas).
- Reducción extrema de texto: Viñetas concisas y directas para máxima limpieza visual.
- Ruteo ortogonal y orden Z-Order optimizado para evitar cualquier superposición de líneas.
- Conectores lineales con cabezas de flecha DrawingML reales.
- Eliminación de bordes oscuros en secciones para lograr un estilo Flat minimalista moderno.
- Guardado robusto con control de PermissionError en caso de bloqueo por PowerPoint abierto.
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def I(v): return Inches(v)

ICONS = r'C:\Servidores\Lectura_FocuzVoz\FocuzVoz3.0\assets\icons'

def ico(name):
    p = os.path.join(ICONS, f'{name}.png')
    return p if os.path.isfile(p) else None

# Paleta de Colores B&W Premium
W   = '#FFFFFF'
BG  = '#FAFAFA'
GRP = '#F4F4F4'
BD  = '#CCCCCC'
BDB = '#333333'
DK  = '#111111'
MD  = '#555555'
SH  = '#E0E0E0'

prs = Presentation()
prs.slide_width  = I(28)
prs.slide_height = I(16)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Fondo principal
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(0), I(0), prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = rgb(BG); bg.line.fill.background()

def section(x, y, w, h, title):
    """Fondo de sección gris suave sin bordes divisorios oscuros (Estilo Flat)."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(GRP)
    s.line.fill.background()
    tb = slide.shapes.add_textbox(I(x+0.15), I(y+0.05), I(w-0.3), I(0.32))
    tf = tb.text_frame; p = tf.paragraphs[0]
    p.text = title; p.alignment = PP_ALIGN.LEFT
    r = p.runs[0]; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = rgb(MD); r.font.name = 'Arial'

def card(x, y, w, h, lib_name, subtitle='', lines=[], icon_key=None):
    """Caja estilizada y limpia con logo a la izquierda, cabecera y viñetas minimalistas."""
    # Sombra
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                I(x+0.04), I(y+0.04), I(w), I(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(SH); s.line.fill.background()
    # Cuerpo
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                I(x), I(y), I(w), I(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(W)
    s.line.color.rgb = rgb(BDB); s.line.width = Pt(1.2)

    # Icono Izquierda
    ico_w = 0.35
    ip = ico(icon_key)
    if ip:
        try:
            slide.shapes.add_picture(ip, I(x+0.1), I(y+0.1), I(ico_w), I(ico_w))
            tx = x + 0.1 + ico_w + 0.08
        except:
            tx = x + 0.1
    else:
        tx = x + 0.1

    # Título (Negrita, 11pt)
    tb = slide.shapes.add_textbox(I(tx), I(y+0.05), I(w - (tx - x) - 0.08), I(0.28))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = lib_name; p.alignment = PP_ALIGN.LEFT
    r = p.runs[0]; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = rgb(DK); r.font.name = 'Arial'

    # Subtítulo (Itálica, 9pt)
    if subtitle:
        tb2 = slide.shapes.add_textbox(I(tx), I(y+0.33), I(w - (tx - x) - 0.08), I(0.20))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]; p2.text = subtitle; p2.alignment = PP_ALIGN.LEFT
        r2 = p2.runs[0]; r2.font.size = Pt(9); r2.font.color.rgb = rgb(MD)
        r2.font.name = 'Arial'; r2.font.italic = True

    # Divisor sutil
    div_y = y + 0.58
    hdiv = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   I(x+0.08), I(div_y), I(w-0.16), I(0.008))
    hdiv.fill.solid(); hdiv.fill.fore_color.rgb = rgb(BD); hdiv.line.fill.background()

    # Viñetas (Letra 9.5pt, espaciadas y con muy poco texto)
    if lines:
        body_y = div_y + 0.08
        tb3 = slide.shapes.add_textbox(I(x+0.1), I(body_y), I(w-0.2), I(h - (body_y - y) - 0.04))
        tf3 = tb3.text_frame; tf3.word_wrap = True
        tf3.margin_left = tf3.margin_right = tf3.margin_top = tf3.margin_bottom = 0
        for i, line in enumerate(lines[:3]):
            p3 = tf3.add_paragraph() if i > 0 else tf3.paragraphs[0]
            p3.text = f"• {line}"
            p3.alignment = PP_ALIGN.LEFT
            r3 = p3.runs[0]; r3.font.size = Pt(9.5); r3.font.color.rgb = rgb(MD)
            r3.font.name = 'Arial'

def lbl(x, y, w, h, text, size=11, bold=False, color=DK,
        align=PP_ALIGN.CENTER, italic=False):
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = rgb(color); r.font.name = 'Arial'

def arr(x1, y1, x2, y2, lw=1.5, arrow_end=True, arrow_start=False):
    """Dibuja un conector lineal directo con flechas opcionales."""
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    c.line.color.rgb = rgb(DK)
    c.line.width = Pt(lw)
    try:
        from pptx.oxml import parse_xml
        line_elem = c.line._get_or_add_ln()
        if arrow_end:
            line_elem.append(parse_xml(
                '<a:tailEnd type="triangle" w="lg" len="lg" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
            ))
        if arrow_start:
            line_elem.append(parse_xml(
                '<a:headEnd type="triangle" w="lg" len="lg" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
            ))
    except Exception as e:
        print(f"Error adding arrowhead: {e}")

def arr_ortho(x1, y1, x2, y2, lw=1.2, x_mid=None, arrow_end=True):
    """Dibuja un conector ortogonal de 3 segmentos (Z-shape) para evitar solapar tarjetas."""
    if x_mid is None:
        x_mid = (x1 + x2) / 2
        
    c1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x_mid), I(y1))
    c2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x_mid), I(y1), I(x_mid), I(y2))
    c3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x_mid), I(y2), I(x2), I(y2))
    
    for c in [c1, c2, c3]:
        c.line.color.rgb = rgb(DK)
        c.line.width = Pt(lw)
        
    if arrow_end:
        try:
            from pptx.oxml import parse_xml
            line_elem = c3.line._get_or_add_ln()
            line_elem.append(parse_xml(
                '<a:tailEnd type="triangle" w="lg" len="lg" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
            ))
        except Exception as e:
            print(f"Error adding arrowhead: {e}")

# ═══════════════════════════════════════════════════════════════════════════════
# TÍTULO PRINCIPAL
lbl(0.2, 0.06, 27.6, 0.5,
    'Arquitectura FocuzVoz 3.0  —  Diagrama de Componentes y Flujo de Datos',
    size=20, bold=True, color=DK)
s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(0.2), I(0.58), I(27.6), I(0.015))
s.fill.solid(); s.fill.fore_color.rgb = rgb(DK); s.line.fill.background()

# ═══════════════════════════════════════════════════════════════════════════════
# 1. RENDERIZADO DE LAS SECCIONES GRISES DE FONDO (Capa inferior)
section(0.2, 0.7, 4.2, 5.5, 'Interfaz de Usuario y Ajustes')
section(0.2, 6.5, 4.2, 8.2, 'Sensores de Entrada')
section(4.8, 0.7, 18.0, 4.5, 'Módulo de Procesamiento de Voz  (voice_controller.py)')
section(4.8, 5.6, 18.0, 4.5, 'Módulo de Procesamiento Visual  (facemesh.py · controllers/)')
section(23.2, 0.7, 4.6, 9.4, 'Sistema Operativo')
section(4.8, 10.5, 23.0, 4.2, 'Persistencia, Telemetría y Distribución')

# ═══════════════════════════════════════════════════════════════════════════════
# 2. RENDERIZADO DE LOS CONECTORES (Capa intermedia)

# Conexiones biológicas desde el usuario (ruteadas limpiamente por la izquierda)
arr_ortho(0.8, 13.1, 0.8, 8.1, x_mid=0.3, lw=1.2) # Voz
arr_ortho(0.8, 13.1, 0.8, 11.1, x_mid=0.5, lw=1.2) # Gesto

# Interfaz GUI ↔ ConfigManager
arr(2.3, 3.1, 2.3, 3.8, lw=1.3, arrow_start=True, arrow_end=True)

# Micrófono (PyAudio) ➔ NumPy (Voice Pipeline Start)
arr_ortho(3.8, 8.1, 6.0, 2.9, x_mid=4.9, lw=1.3)

# Flujo horizontal del Pipeline de Voz (NumPy -> Vosk -> pyttsx3 -> pynput)
arr(9.0, 2.9, 10.2, 2.9, lw=1.3)
arr(13.2, 2.9, 14.4, 2.9, lw=1.3)
arr(17.4, 2.9, 18.6, 2.9, lw=1.3)

# Cámara (OpenCV) ➔ MediaPipe (Vision Pipeline Start)
arr_ortho(3.8, 10.6, 6.0, 7.8, x_mid=5.3, lw=1.3)

# Flujo horizontal del Pipeline de Visión (MediaPipe -> OneEuro -> FSM -> Mouse)
arr(9.0, 7.8, 10.2, 7.8, lw=1.3)
arr(13.2, 7.8, 14.4, 7.8, lw=1.3)
arr(17.4, 7.8, 18.6, 7.8, lw=1.3)

# Conexiones hacia la capa de ejecución del Sistema Operativo
arr(21.6, 2.9, 24.0, 3.5, lw=1.5)
arr(21.6, 7.8, 24.0, 7.2, lw=1.5)

# PyDirectInput to Windows OS
arr(25.5, 4.4, 25.5, 6.3, lw=1.3)

# Conector ConfigManager ➔ SQLite3
arr_ortho(3.8, 4.7, 6.95, 12.6, x_mid=4.1, lw=1.3)

# ConfigManager propaga configuraciones dinámicas de perfiles
arr_ortho(3.8, 4.2, 6.0, 2.4, x_mid=4.5, lw=1.1)
arr_ortho(3.8, 4.2, 6.0, 7.2, x_mid=4.5, lw=1.1)

# GUI controla la previsualización directa de cámara (ruteada por la izquierda a diferente entrada)
arr_ortho(0.8, 2.2, 0.8, 10.1, x_mid=0.15, lw=1.1)

# Conexión vertical de telemetría (logs) hacia Research Analytics
arr(15.9, 3.8, 15.9, 11.7, lw=1.1)
arr(15.9, 8.7, 15.9, 11.7, lw=1.1)

# Persistencia de telemetría a la base de datos local
arr(13.9, 12.6, 11.75, 12.6, lw=1.1)

# Proceso de compilación EXE nativo
arr_ortho(25.5, 8.1, 23.25, 11.7, x_mid=25.5, lw=1.1)

# ═══════════════════════════════════════════════════════════════════════════════
# 2.5 RENDERIZADO DE LAS ETIQUETAS DE TEXTO (Manual y sin encimes)

# Ajustes
lbl(2.38, 3.35, 1.5, 0.25, 'Ajustes GUI', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)

# Entradas desde Usuario (Izquierda)
lbl(0.35, 11.8, 1.2, 0.25, 'Comando de voz', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)
lbl(0.55, 12.1, 1.2, 0.25, 'Gesto facial', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)
lbl(0.2, 6.0, 1.5, 0.25, 'Preview de cámara', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)

# Entrada a los Pipelines
lbl(5.0, 2.6, 1.5, 0.25, 'Audio PCM (16kHz)', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)
lbl(4.0, 10.3, 1.5, 0.25, 'Video RGB (30 FPS)', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)

# Propagación de Configuración
lbl(4.6, 2.1, 1.3, 0.25, 'Umbrales RMS', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)
lbl(4.6, 7.0, 1.3, 0.25, 'Filtro cursor', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)
lbl(4.3, 12.3, 2.5, 0.25, 'Cargar perfiles (JSON/DB)', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)

# Pipeline de Voz (Internos)
lbl(9.1, 2.6, 1.0, 0.25, 'Audio filtrado', size=8.5, bold=True, color=MD, align=PP_ALIGN.CENTER)
lbl(13.25, 2.6, 1.1, 0.25, 'Texto reconocido', size=8.5, bold=True, color=MD, align=PP_ALIGN.CENTER)
lbl(17.45, 2.6, 1.1, 0.25, 'Confirmación TTS', size=8.5, bold=True, color=MD, align=PP_ALIGN.CENTER)

# Pipeline de Visión (Internos)
lbl(9.1, 7.5, 1.0, 0.25, 'Landmarks 3D', size=8.5, bold=True, color=MD, align=PP_ALIGN.CENTER)
lbl(13.25, 7.5, 1.1, 0.25, 'Coordenadas X/Y', size=8.5, bold=True, color=MD, align=PP_ALIGN.CENTER)
lbl(17.45, 7.5, 1.1, 0.25, 'Gesto / Dwell', size=8.5, bold=True, color=MD, align=PP_ALIGN.CENTER)

# Capa de Ejecución y OS
lbl(21.7, 2.6, 1.5, 0.25, 'Simulación teclas', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)
lbl(21.7, 7.5, 1.5, 0.25, 'Movimiento mouse', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)
lbl(25.6, 5.2, 1.4, 0.25, 'Inyección Win32', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)
lbl(23.3, 11.4, 1.5, 0.25, 'Compilar EXE', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)

# Telemetría y Persistencia
lbl(16.0, 5.0, 1.0, 0.25, 'Logs STT', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)
lbl(16.0, 9.5, 1.0, 0.25, 'Logs Gestos', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)
lbl(11.9, 12.3, 1.5, 0.25, 'Persistir logs', size=8.5, bold=True, color=MD, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
# 3. RENDERIZADO DE LAS TARJETAS (Capa superior)

# Control e Interfaz (Columna izquierda)
card(0.8, 1.3, 3.0, 1.8, 'CustomTkinter GUI', 'views/ vistas del panel', [
    'Interfaz de control principal',
    'Configuración visual de perfiles',
    'Previsualización en tiempo real'
], icon_key='python')

card(0.8, 3.8, 3.0, 1.8, 'ConfigManager', 'Singleton de Ajustes', [
    'Persistencia JSON y SQLite',
    'Sincronización de perfiles',
    'Actualización en caliente'
], icon_key='python')

# Sensores de Entrada (Columna izquierda)
card(0.8, 7.2, 3.0, 1.8, 'PyAudio', 'Captura de micrófono', [
    'Audio PCM mono a 16000 Hz',
    'Búfer de entrada dinámico',
    'Ejecución en hilo secundario'
], icon_key='pyaudio')

card(0.8, 9.7, 3.0, 1.8, 'OpenCV', 'camera_manager.py', [
    'Captura a 30 FPS en 4:3',
    'Filtro espejo y recorte',
    'Hilo de cámara independiente'
], icon_key='opencv')

card(0.8, 12.2, 3.0, 1.8, 'Usuario', 'Entrada biológica', [
    'Movimiento y gestos faciales',
    'Comandos de voz directos',
    'Operación manos libres'
], icon_key=None)

# Pipeline de Voz (Fila superior - NumPy -> Vosk -> pyttsx3 -> pynput)
card(6.0, 2.0, 3.0, 1.8, 'NumPy', 'Filtro Noise Gate RMS', [
    'Supresión activa de ruido',
    'Cálculo de volumen RMS',
    'Umbral de volumen dinámico'
], icon_key='numpy')

card(10.2, 2.0, 3.0, 1.8, 'Vosk STT', 'Reconocimiento offline', [
    'Modelo local en español',
    'Transcripción offline rápida',
    'Procesamiento en memoria'
], icon_key='vosk')

card(14.4, 2.0, 3.0, 1.8, 'pyttsx3', 'Feedback Auditivo', [
    'Confirmación de voz',
    'Hilo dedicado COM seguro',
    'Síntesis de voz offline'
], icon_key='python')

card(18.6, 2.0, 3.0, 1.8, 'pynput / pyautogui', 'Traductor de Comandos', [
    'Conversión de texto a acción',
    'Cierre seguro (TaskKiller)',
    'Simulación de teclas físicas'
], icon_key='python')

# Pipeline de Visión (Fila inferior)
card(6.0, 6.9, 3.0, 1.8, 'MediaPipe', 'FaceLandmarker', [
    'Detección 478 landmarks',
    '52 blendshapes faciales',
    'Matriz de rotación 3D'
], icon_key='mediapipe')

card(10.2, 6.9, 3.0, 1.8, 'OneEuroFilter2D', 'Filtro de Ruido', [
    'Filtro adaptativo de cursor',
    'Eliminación total de jitter',
    'Suavizado en tiempo real'
], icon_key='python')

card(14.4, 6.9, 3.0, 1.8, 'FacialEventManager', 'Máquina de Estados FSM', [
    'Control Midas Touch (150ms)',
    'Gestión de gestos y dwell',
    'Mapeo de acciones y delay'
], icon_key='python')

card(18.6, 6.9, 3.0, 1.8, 'MouseController', 'Gestor de Puntero', [
    'Movimiento suave por rostro',
    'Aceleración sigmoidal',
    'Soporte multimonitor nativo'
], icon_key='python')

# Capa de Ejecución OS (Columna derecha)
card(24.0, 2.6, 3.0, 1.8, 'PyDirectInput / pyautogui', 'Inyector de Acciones', [
    'Inyección a bajo nivel',
    'Compatibilidad con juegos',
    'Simulación de pulsaciones'
], icon_key='python')

card(24.0, 6.3, 3.0, 1.8, 'Windows OS', 'Recepción de Eventos', [
    'API user32.dll nativa',
    'Recepción nativa de eventos',
    'Simulación vía Win32'
], icon_key='windows')

# Persistencia y Empaquetado (Fila inferior)
card(6.95, 11.7, 4.8, 1.8, 'SQLite3 DatabaseManager', 'focuzvoz.db local', [
    'Almacenamiento de perfiles',
    'Base de datos relacional',
    'Integridad por llaves foráneas'
], icon_key='sqlite')

card(13.9, 11.7, 4.8, 1.8, 'Research Analytics', 'Telemetría Experimental', [
    'Registro de telemetría local',
    'Exportación de logs en CSV',
    'Métricas de clics y comandos'
], icon_key='python')

card(20.85, 11.7, 4.8, 1.8, 'PyInstaller / Inno Setup', 'Empaquetado y Distribución', [
    'Compilación a ejecutable único',
    'Empaquetado de assets locales',
    'Instalador offline (.exe)'
], icon_key='pyinstaller')

# ═══════════════════════════════════════════════════════════════════════════════
# Leyenda inferior explicativa de diseño e implementación
lbl(0.2, 15.2, 27.6, 0.28,
    'Monolítico Modular  ·  Singletons (Metaclase)  ·  '
    'Threading + ThreadPoolExecutor  ·  '
    'Comunicación por Memoria Compartida (numpy ndarray)  ·  '
    'Latencia < 16ms  ·  100% Offline',
    size=9.5, italic=True, color=MD)

# Guardado seguro en caso de archivos de PowerPoint abiertos por el usuario
out = 'Arquitectura_FocuzVoz_Final_v17.pptx'
try:
    prs.save(out)
    print(f'[OK]  {out}')
except PermissionError:
    # Si está abierto, intentamos guardar con un sufijo incremental
    saved = False
    for idx in range(16, 30):
        alt_out = f'Arquitectura_FocuzVoz_Final_v{idx}.pptx'
        try:
            prs.save(alt_out)
            print(f'[OK]  {alt_out} (El archivo principal estaba abierto en PowerPoint, se guardó copia alternativa)')
            saved = True
            break
        except PermissionError:
            continue
    if not saved:
        print('[ERROR] No se pudo guardar el archivo. Por favor, cierra PowerPoint y reintenta.')
        sys.exit(1)
