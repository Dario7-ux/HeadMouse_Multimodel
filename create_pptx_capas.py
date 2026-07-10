import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def U(v): return Inches(v / 100.0)

# ==============================================================================
# CONFIGURACIÓN DE INFOGRAFÍA (Vertical, 10x14 pulgadas)
# ==============================================================================
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(14)
BLANK = prs.slide_layouts[6]

C = {
    'bg_dark':   '#FFFFFF', 'bg_card':   '#F8F9FA', 'accent1':   '#2B2B2B',
    'accent2':   '#4A4A4A', 'accent3':   '#6C757D', 'text_head': '#111111',
    'text_body': '#333333', 'text_dim':  '#555555', 'border':    '#DEE2E6',
    'layer_in':  '#E9ECEF', 'green_lt':  '#000000', 'blue_lt':   '#222222',
    'orange_lt': '#333333', 'purple_lt': '#444444', 'red_lt':    '#555555',
}

def make_helpers(slide):
    def bg(color=C['bg_dark']):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        sh.fill.solid(); sh.fill.fore_color.rgb = hex_rgb(color); sh.line.fill.background()
    def rect(x, y, w, h, fill, stroke=None, radius=True):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, U(x), U(y), U(w), U(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = hex_rgb(fill)
        if stroke: sh.line.color.rgb = hex_rgb(stroke); sh.line.width = Pt(2.0)
        else: sh.line.fill.background()
        return sh
    def label(x, y, w, h, text, size=14, bold=False, color=C['text_body'], align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(U(x), U(y), U(w), U(h))
        tf = tb.text_frame; tf.word_wrap = wrap; p = tf.paragraphs[0]; p.text = text; p.alignment = align
        run = p.runs[0]; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = hex_rgb(color); run.font.name = 'Segoe UI'
        return tb
    def arrow(x1, y1, x2, y2, color=C['text_dim'], width=2.5):
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, U(x1), U(y1), U(x2), U(y2))
        conn.line.color.rgb = hex_rgb(color); conn.line.width = Pt(width)
    def chip(x, y, w, h, text, fill, text_color=C['text_head'], size=12, bold=False):
        rect(x, y, w, h, fill, stroke=None, radius=True)
        label(x, y + (h/2 - 12) - (size-12)*2, w, h, text, size=size, bold=bold, color=text_color, align=PP_ALIGN.CENTER)
    return bg, rect, label, arrow, chip


# ==============================================================================
#  SLIDE ÚNICA: LA INFOGRAFÍA DE ARQUITECTURA
# ==============================================================================
s1 = prs.slides.add_slide(BLANK)
bg, rect, label, arrow, chip = make_helpers(s1)
bg()

# --- HEADER ---
rect(0, 0, 1000, 10, C['accent1'])
label(50, 40, 900, 60, 'FocuzVoz 3.0', size=42, bold=True, color=C['text_head'], align=PP_ALIGN.CENTER)
label(50, 100, 900, 40, 'Arquitectura de Software (Módulos y Flujo de Datos)', size=18, color=C['text_dim'], align=PP_ALIGN.CENTER)
rect(350, 150, 300, 2, C['border'])

# ==============================================================================
#  ARQUITECTURA POR CAPAS (Y = 180 a 1150)
# ==============================================================================

# --- CAPA 1: PRESENTACIÓN (UI) ---
rect(50, 200, 900, 220, C['layer_in'], stroke=C['border'])
label(50, 220, 900, 30, 'CAPA DE PRESENTACIÓN (Interfaz y Estado)', size=16, bold=True, color=C['text_head'], align=PP_ALIGN.CENTER)

chip(80, 280, 250, 110, '🖥️ Interfaz Usuario\nsrc/gui/main_gui.py\nPanel de Configuración', C['bg_card'], text_color=C['text_body'], size=12)
chip(375, 280, 250, 110, '🕹️ Overlay UI\nsrc/gui/bot_overlay.py\nFeedback Visual en Pantalla', C['bg_card'], text_color=C['text_body'], size=12)
chip(670, 280, 250, 110, '⚙️ Gestor de Perfiles\nsrc/config_manager.py\nEstado Global (Singleton)', C['bg_card'], text_color=C['text_body'], size=12)

# Flechas entre Capa 1 y Capa 2
arrow(330, 420, 330, 500, C['text_head'], 3)
label(130, 440, 180, 40, 'Actualiza\nConfiguraciones', size=11, align=PP_ALIGN.RIGHT)

arrow(670, 500, 670, 420, C['text_head'], 3)
label(690, 440, 180, 40, 'Emite Feedback\nal Usuario', size=11, align=PP_ALIGN.LEFT)


# --- CAPA 2: LÓGICA DE NEGOCIO ---
rect(50, 500, 900, 360, C['layer_in'], stroke=C['border'])
label(50, 520, 900, 30, 'CAPA DE LÓGICA DE NEGOCIO (Motor Central y Controladores)', size=16, bold=True, color=C['text_head'], align=PP_ALIGN.CENTER)

chip(80, 580, 250, 120, '👁️ Pipeline de Visión\nsrc/camera_manager.py\nsrc/detectors/facemesh.py\nMediaPipe Blendshapes', C['bg_card'], text_color=C['text_body'], size=12)
chip(375, 580, 250, 120, '🧠 Gestor de Eventos\nsrc/controllers/facial_event_manager.py\nDwell Time y Cooldown', C['bg_card'], text_color=C['text_body'], size=12)
chip(670, 580, 250, 120, '🎙️ Procesador de Voz\nsrc/controllers/voice_controller.py\nVosk Offline STT', C['bg_card'], text_color=C['text_body'], size=12)

chip(200, 740, 600, 80, '🔄 Filtrado y Orquestación\nsrc/utils/one_euro_filter.py (Suavizado) | src/pipeline.py', C['bg_card'], text_color=C['text_body'], size=12)

# Flechas entre Capa 2 y Capa 3
arrow(200, 980, 200, 860, C['accent1'], 3)
label(40, 900, 150, 40, 'Frames Video\ny Audio PCM', size=11, bold=True, align=PP_ALIGN.RIGHT)

arrow(500, 860, 500, 980, C['text_head'], 3)
label(520, 900, 180, 40, 'Consultas y Logs\nhacia focuzvoz.db', size=11)

arrow(800, 860, 800, 980, C['accent1'], 3)
label(820, 900, 150, 40, 'Comandos al SO\n(Clicks y Teclas)', size=11, bold=True)


# --- CAPA 3: INFRAESTRUCTURA ---
rect(50, 980, 900, 220, C['layer_in'], stroke=C['border'])
label(50, 1000, 900, 30, 'CAPA DE INFRAESTRUCTURA (Hardware, SO y Datos)', size=16, bold=True, color=C['text_head'], align=PP_ALIGN.CENTER)

chip(80, 1060, 250, 110, '📥 Entradas de Hardware\nOpenCV (Cámara Web)\nPyAudio (Micrófono)', C['bg_card'], text_color=C['text_body'], size=12)
chip(375, 1060, 250, 110, '💾 Persistencia\nsrc/utils/database.py\nBase de Datos SQLite3', C['bg_card'], text_color=C['text_body'], size=12)
chip(670, 1060, 250, 110, '📤 Salidas de Sistema\nmouse_controller.py\nkeybinder.py (PyAutoGUI)', C['bg_card'], text_color=C['text_body'], size=12)


# --- PIE DE PÁGINA ---
rect(0, 1340, 1000, 60, C['layer_in'])
label(50, 1360, 900, 30, 'FocuzVoz 3.0 · Infografía de Arquitectura Monolítica por Capas · Comunicación mediante callbacks e hilos paralelos.', size=12, color=C['text_dim'], align=PP_ALIGN.CENTER)

output = 'Infografia_Arquitectura_Final.pptx'
prs.save(output)
print(f'[OK]  Infografia de arquitectura generada: {output}')
