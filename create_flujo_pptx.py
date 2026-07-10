
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def U(v):
    return Inches(v / 100.0)

prs = Presentation()
prs.slide_width  = Inches(18)
prs.slide_height = Inches(10.125)
BLANK = prs.slide_layouts[6]

C = {
    'bg':        '#FFFFFF',
    'bg_card':   '#F8F9FA',
    'layer_in':  '#E9ECEF',
    'border':    '#CED4DA',
    'head':      '#000000',
    'body':      '#212529',
    'dim':       '#6C757D',
    'black':     '#000000',
    'dark1':     '#212529',
    'dark2':     '#343A40',
    'dark3':     '#495057',
    'gray':      '#6C757D',
    'light1':    '#ADB5BD',
    'light2':    '#CED4DA',
    'light3':    '#E9ECEF',
    'light4':    '#F8F9FA',
    'white':     '#FFFFFF',
}

def make_helpers(slide):
    def bg(color=C['bg']):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        sh.fill.solid(); sh.fill.fore_color.rgb = hex_rgb(color)
        sh.line.fill.background()

    def rect(x, y, w, h, fill, stroke=None, radius=False):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        sh = slide.shapes.add_shape(kind, U(x), U(y), U(w), U(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = hex_rgb(fill)
        if stroke:
            sh.line.color.rgb = hex_rgb(stroke)
            sh.line.width = Pt(2.0)
        else:
            sh.line.fill.background()
        return sh

    def diamond(x, y, w, h, fill, stroke=None):
        sh = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, U(x), U(y), U(w), U(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = hex_rgb(fill)
        if stroke:
            sh.line.color.rgb = hex_rgb(stroke)
            sh.line.width = Pt(2.0)
        else:
            sh.line.fill.background()
        return sh

    def oval(x, y, w, h, fill, stroke=None):
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, U(x), U(y), U(w), U(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = hex_rgb(fill)
        if stroke:
            sh.line.color.rgb = hex_rgb(stroke)
            sh.line.width = Pt(2.0)
        else:
            sh.line.fill.background()
        return sh

    def label(x, y, w, h, text, size=13, bold=False, color=C['body'],
              align=PP_ALIGN.CENTER, wrap=True, font='Segoe UI'):
        tb = slide.shapes.add_textbox(U(x), U(y), U(w), U(h))
        tf = tb.text_frame; tf.word_wrap = wrap
        p = tf.paragraphs[0]; p.text = text; p.alignment = align
        run = p.runs[0]
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = hex_rgb(color); run.font.name = font
        return tb

    def arrow(x1, y1, x2, y2, color=C['dark3'], width=2.5):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, U(x1), U(y1), U(x2), U(y2))
        conn.line.color.rgb = hex_rgb(color)
        conn.line.width = Pt(width)

    def badge(x, y, w, h, text, fill, tcol, size=12, bold=True):
        rect(x, y, w, h, fill, stroke=None, radius=True)
        label(x, y, w, h, text, size=size, bold=bold,
              color=tcol, align=PP_ALIGN.CENTER)

    return bg, rect, diamond, oval, label, arrow, badge

# ══════════════════════════════════════════════════════════════
#  SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg, rect, diamond, oval, label, arrow, badge = make_helpers(s1)
bg()
rect(0, 0, 1800, 8, C['dark1'])
label(80, 270, 1640, 80,
      'FocuzVoz 3.0 — Diagrama de Flujo del Sistema',
      size=44, bold=True, color=C['head'], align=PP_ALIGN.CENTER)
label(80, 370, 1640, 50,
      'Procesamiento por Frame · Control Facial · Dictado por Voz · Clics Oculares',
      size=22, color=C['dim'], align=PP_ALIGN.CENTER)

badge(660, 460, 480, 56,
      '● Versión 3.0  |  Aplicación Desktop Python',
      C['layer_in'], C['dark1'], size=16)

rect(0, 992, 1800, 4, C['dark2'])
label(80, 960, 1640, 28,
      'Procesamiento facial offline · Latencia < 16 ms por frame · Privacidad total sin conexión',
      size=14, color=C['dim'], align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  SLIDE 2 — PARÁMETROS CONFIGURABLES
# ══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg, rect, diamond, oval, label, arrow, badge = make_helpers(s2)
bg()
rect(0, 0, 1800, 8, C['dark1'])
label(80, 14, 1640, 44,
      'Parámetros Configurables del Sistema — Valores de Referencia',
      size=30, bold=True, color=C['head'], align=PP_ALIGN.CENTER)

params = [
    (C['dark2'], C['head'], '👁  Tiempo Mínimo de Guiño',
     'Duración mínima del parpadeo para ser\nconsiderado guiño intencional.',
     '150 – 300 ms', 'Valor inicial: 200 ms\nRango configurable: 100 – 400 ms'),
    (C['dark2'], C['head'], '⏱  Dwell Time (Permanencia)',
     'Tiempo que el gesto debe mantenerse\nantes de activar una acción.',
     '100 – 800 ms', 'Valor inicial: 150 ms\nRango configurable: 100 ms – 1.0 s'),
    (C['dark2'], C['head'], '🔁  Cooldown (Enfriamiento)',
     'Período de bloqueo tras ejecutar\nuna acción para evitar repetición.',
     '200 – 1000 ms', 'Valor inicial: 350 ms\nRango configurable: 200 ms – 2.0 s'),
    (C['dark2'], C['head'], '📊  Umbral de Blendshape',
     'Valor mínimo del coeficiente de blendshape\npara reconocer el gesto como activo.',
     '0.35 – 0.55', 'Valor inicial: 0.40\nRango: 0.20 (sensible) – 0.70 (estricto)'),
    (C['dark2'], C['head'], '🔊  Umbral RMS de Audio',
     'Amplitud mínima de señal para activar\nel pipeline de reconocimiento de voz.',
     '100 – 1000', 'Valor inicial: 500 (Sens. 50)\nRango configurable: 0 (100) – 1000 (0)'),
    (C['dark2'], C['head'], '🎯  Sensibilidad del Cursor',
     'Velocidad de desplazamiento base\nhacia coordenadas de pantalla.',
     '10.0 – 30.0', 'Valor inicial: X=18.0, Y=22.0\nFiltro One Euro: β=0.015, f_c=2.5 Hz'),
]

card_w, card_h = 540, 350
cols = [40, 620, 1200]
rows = [75, 455]
positions = [(c, r) for r in rows for c in cols]

for i, (stroke, tcol, title, desc, rng, detail) in enumerate(params):
    cx, cy = positions[i]
    rect(cx, cy, card_w, card_h, C['bg_card'], stroke=stroke, radius=True)
    rect(cx, cy, card_w, 48, C['light2'], stroke=None, radius=False)
    label(cx+14, cy+8, card_w-28, 36, title,
          size=16, bold=True, color=tcol, align=PP_ALIGN.LEFT)
    label(cx+18, cy+60, card_w-36, 76, desc,
          size=14, color=C['body'], align=PP_ALIGN.LEFT)
    badge(cx+18, cy+146, card_w-36, 50,
          f'Rango:  {rng}', C['light3'], tcol, size=18, bold=True)
    label(cx+18, cy+210, card_w-36, 130, detail,
          size=13, color=C['dim'], align=PP_ALIGN.LEFT)

rect(0, 984, 1800, 4, C['dark2'])
label(80, 954, 1640, 26,
      'Todos los parámetros son ajustables en tiempo real desde el panel GUI · configs/focuzvoz.db · Perfiles por usuario',
      size=13, color=C['dim'], align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  SLIDE 3 — DIAGRAMA DE FLUJO PRINCIPAL (Grayscale & Spaced)
# ══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg, rect, diamond, oval, label, arrow, badge = make_helpers(s3)
bg()
rect(0, 0, 1800, 8, C['dark1'])

label(80, 12, 1640, 40,
      'FocuzVoz 3.0 — Diagrama de Flujo por Frame de Procesamiento',
      size=26, bold=True, color=C['head'], align=PP_ALIGN.CENTER)

CX = 900
Y0 = 60
oval(CX-110, Y0, 220, 56, C['white'], stroke=C['dark1'])
label(CX-110, Y0+2, 220, 52, 'START', size=20, bold=True, color=C['dark1'])

Y1 = 146
rect(CX-150, Y1, 300, 60, C['bg_card'], stroke=C['dark2'], radius=True)
label(CX-150, Y1, 300, 60,
      'Activar Webcam\ny Micrófono', size=14, bold=True, color=C['dark2'])

Y2 = 236
oval(CX-110, Y2, 220, 56, C['light3'], stroke=C['dim'])
label(CX-110, Y2+2, 220, 52, 'Loop Start', size=15, bold=True, color=C['dark2'])

Y3 = 322
rect(CX-160, Y3, 320, 60, C['bg_card'], stroke=C['dark2'], radius=True)
label(CX-160, Y3, 320, 60,
      'Capturar Frame\ny Bloque de Audio', size=14, bold=True, color=C['dark2'])

Y4 = 412
diamond(CX-140, Y4, 280, 80, C['light4'], stroke=C['dark3'])
label(CX-140, Y4, 280, 80,
      '¿Malla facial\nválida?', size=14, bold=True, color=C['dark3'])

Y5 = 522
rect(CX-180, Y5, 360, 60, C['bg_card'], stroke=C['dark1'], radius=True)
label(CX-180, Y5, 360, 60,
      'Extraer coords 3D\nnose, eyes, mouth', size=14, bold=True, color=C['dark1'])

# ─────────────────────────────────────────────────────────
#  RAMAS
# ─────────────────────────────────────────────────────────
YBASE = 612

# ── Branch 1
X1 = 20
rect(X1, YBASE, 380, 280, C['bg_card'], stroke=C['dark2'], radius=False)
label(X1, YBASE+8, 380, 26, 'Rama 1 · Control de Cursor', size=13, bold=True, color=C['dark2'])

badge(X1+16, YBASE+46, 348, 42, 'Filtro One Euro → Posición nariz', C['light3'], C['dark2'], size=13)
badge(X1+16, YBASE+100, 348, 42, 'Mapear a resolución de pantalla', C['light3'], C['dark2'], size=13)
badge(X1+16, YBASE+154, 348, 42, 'Aplicar aceleración sigmoide', C['light3'], C['dark2'], size=13)
badge(X1+16, YBASE+208, 348, 42, 'Mover puntero en OS (PyAutoGUI)', C['light3'], C['dark2'], size=13)
label(X1, YBASE+252, 380, 24, 'Sens (X/Y): 18.0 / 22.0  |  f_c=2.5 Hz  |  β=0.015', size=11, color=C['dim'])

# ── Branch 2
X2 = 420
rect(X2, YBASE, 360, 280, C['bg_card'], stroke=C['dark2'], radius=False)
label(X2, YBASE+8, 360, 26, 'Rama 2 · Clics Oculares', size=13, bold=True, color=C['dark2'])

diamond(X2+50, YBASE+42, 260, 80, C['light4'], stroke=C['dark3'])
label(X2+50, YBASE+42, 260, 80, '¿Guiño\ndetectado?', size=14, bold=True, color=C['dark3'])

label(X2, YBASE+126, 360, 20, 'Blendshape ≥ 0.40', size=11, color=C['dim'])
label(X2, YBASE+144, 360, 20, 'Duración ≥ 150 ms', size=11, color=C['dim'])
label(X2, YBASE+162, 360, 20, 'Cooldown: 350 ms', size=11, color=C['dim'])

badge(X2+16, YBASE+186, 328, 44, 'Ejecutar Clic OS (PyDirectInput)', C['light3'], C['dark2'], size=13)
label(X2, YBASE+234, 360, 40, 'SÍ: Clic · NO: Sin acción\nDwell Time: 150 ms', size=12, bold=True, color=C['dark3'])

label(X2+300, YBASE+94, 40, 20, 'SÍ', size=12, bold=True, color=C['dark1'])
label(X2+14,  YBASE+94, 40, 20, 'NO', size=12, bold=True, color=C['dark1'])

# ── Branch 3
X3 = 800
rect(X3, YBASE, 980, 280, C['bg_card'], stroke=C['dark2'], radius=False)
label(X3, YBASE+8, 980, 26, 'Rama 3 · Comandos y Dictado por Voz', size=13, bold=True, color=C['dark2'])

XVL = X3 + 30
diamond(XVL, YBASE+42, 260, 76, C['light4'], stroke=C['dark3'])
label(XVL, YBASE+42, 260, 76, '¿Boca\nabierta?', size=14, bold=True, color=C['dark3'])
badge(XVL+10, YBASE+136, 240, 44, 'Filtro RMS / Noise Gate', C['light3'], C['dark2'], size=13)
diamond(XVL, YBASE+196, 260, 70, C['light4'], stroke=C['dark3'])
label(XVL, YBASE+196, 260, 70, '¿RMS ≥ 0.04?', size=14, bold=True, color=C['dark3'])
# badge(XVL+10, YBASE+240, 240, 40, 'Decodificar (Vosk)', C['light3'], C['dark2'], size=12) # Se omite por espacio, o lo hacemos más ajustado.
# Mejor ajusto alturas:
badge(XVL+10, YBASE+128, 240, 36, 'Filtro RMS / Noise Gate', C['light3'], C['dark2'], size=12)
diamond(XVL, YBASE+176, 260, 60, C['light4'], stroke=C['dark3'])
label(XVL, YBASE+176, 260, 60, '¿RMS ≥ 500?', size=13, bold=True, color=C['dark3'])
badge(XVL+10, YBASE+242, 240, 36, 'Decodificar (Vosk)', C['light3'], C['dark2'], size=12)

XVR = X3 + 360
diamond(XVR, YBASE+42, 280, 76, C['light4'], stroke=C['dark3'])
label(XVR, YBASE+42, 280, 76, '¿Comando\no texto?', size=14, bold=True, color=C['dark3'])
badge(XVR+10, YBASE+136, 260, 44, 'Ejecutar atajo / Escribir', C['light3'], C['dark2'], size=13)

label(XVR, YBASE+190, 280, 24, 'Wake-word: "focuz"', size=12, color=C['dim'])
label(XVR, YBASE+214, 280, 24, 'Umbral RMS: 500 | Noise Gate ON', size=12, color=C['dim'])
label(XVR, YBASE+238, 280, 24, 'SÍ: Acción · NO: ignorar', size=12, bold=True, color=C['dark3'])

label(XVL-32, YBASE+70, 36, 20, 'NO', size=12, bold=True, color=C['dark1'])
label(XVL+250, YBASE+70, 36, 20, 'SÍ', size=12, bold=True, color=C['dark1'])
label(XVL-32, YBASE+196, 36, 20, 'NO', size=12, bold=True, color=C['dark1'])
label(XVL+250, YBASE+196, 36, 20, 'SÍ', size=12, bold=True, color=C['dark1'])

label(XVR-32, YBASE+70, 36, 20, 'NO', size=12, bold=True, color=C['dark1'])
label(XVR+270, YBASE+70, 36, 20, 'SÍ', size=12, bold=True, color=C['dark1'])

# ─────────────────────────────────────────────────────────
#  End of Frame
# ─────────────────────────────────────────────────────────
YEND = 910
oval(CX-110, YEND, 220, 56, C['white'], stroke=C['dark1'])
label(CX-110, YEND+2, 220, 52, 'End of Frame', size=16, bold=True, color=C['dark1'])

# ─────────────────────────────────────────────────────────
#  FLECHAS
# ─────────────────────────────────────────────────────────
MX = CX

arrow(MX, Y0+56, MX, Y1, C['dark3'], 2.5)
arrow(MX, Y1+60, MX, Y2, C['dark3'], 2.5)
arrow(MX, Y2+56, MX, Y3, C['dark3'], 2.5)
arrow(MX, Y3+60, MX, Y4, C['dark3'], 2.5)
arrow(MX, Y4+80, MX, Y5, C['dark3'], 2.5)

# Malla NO -> loop
arrow(MX+140, Y4+40, 1750, Y4+40, C['dark2'], 2.5)
arrow(1750, Y4+40, 1750, Y2+28, C['dark2'], 2.5)
arrow(1750, Y2+28, MX+110, Y2+28, C['dark2'], 2.5)
label(1710, Y4+14, 40, 24, 'NO', size=13, bold=True, color=C['dark1'])

arrow(MX, Y5+60, MX, YBASE-16, C['dark3'], 2.5)
arrow(210, YBASE-16, 1600, YBASE-16, C['dark3'], 2.0)
arrow(210, YBASE-16, 210, YBASE, C['dark3'], 2.5)
arrow(600, YBASE-16, 600, YBASE, C['dark3'], 2.5)
arrow(1000, YBASE-16, 1000, YBASE, C['dark3'], 2.5)

YBOT = YBASE + 280
arrow(210, YBOT, 210, YEND+28, C['dark3'], 2.5)
arrow(600, YBOT, 600, YEND+28, C['dark3'], 2.5)
arrow(1000, YBOT, 1000, YEND+28, C['dark3'], 2.5)
arrow(210, YEND+28, MX-110, YEND+28, C['dark3'], 2.0)
arrow(MX+110, YEND+28, 1000, YEND+28, C['dark3'], 2.0)

# Loop retorno
XLFT = 12
arrow(MX-110, YEND+28, XLFT, YEND+28, C['dark2'], 2.5)
arrow(XLFT, YEND+28, XLFT, Y2+28, C['dark2'], 2.5)
arrow(XLFT, Y2+28, MX-110, Y2+28, C['dark2'], 2.5)
label(20, Y2-10, 120, 24, 'Next Iteration', size=13, bold=True, color=C['dark1'], align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════
output = 'Flujo_FocuzVoz_Ramas.pptx'
prs.save(output)
print(f'[OK]  Presentacion generada: {output}  (3 diapositivas)')
